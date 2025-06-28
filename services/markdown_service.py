from pptx import Presentation
from PIL import Image
import os
from io import BytesIO
import easyocr
import numpy as np

def pptx_to_markdown_text(pptx_path, lang="fr"):
    prs = Presentation(pptx_path)
    filename = os.path.splitext(os.path.basename(pptx_path))[0]
    output_dir = "images_extraites"
    os.makedirs(output_dir, exist_ok=True)

    # Initialisation de l'OCR EasyOCR
    reader = easyocr.Reader([lang], gpu=False)

    md_lines = [f"# Présentation : {filename}\n"]

    for i, slide in enumerate(prs.slides):
        md_lines.append(f"\n---\n\n## Slide {i + 1}\n")
        image_counter = 1

        for shape in slide.shapes:
            # Cas 1 : Texte (zones de texte, formes, etc.)
            if hasattr(shape, "text") and shape.text.strip():
                md_lines.append(shape.text.strip().replace('\r', '') + "\n")

            # Cas 2 : Tableaux
            elif shape.has_table:
                table = shape.table
                headers = [cell.text.strip() for cell in table.rows[0].cells]
                md_lines.append("| " + " | ".join(headers) + " |")
                md_lines.append("|" + "|".join([" --- " for _ in headers]) + "|")

                for row_idx in range(1, len(table.rows)):
                    row = table.rows[row_idx]
                    cells = [cell.text.strip() for cell in row.cells]
                    md_lines.append("| " + " | ".join(cells) + " |")

            # Cas 3 : Images avec OCR
            elif shape.shape_type == 13 and hasattr(shape, "image"):
                try:
                    image = shape.image
                    image_bytes = image.blob
                    pil_image = Image.open(BytesIO(image_bytes))

                    # Sauvegarde optionnelle de l’image
                    image_filename = f"slide{i+1}_img{image_counter}.png"
                    image_path = os.path.join(output_dir, image_filename)
                    pil_image.save(image_path)

                    # Conversion PIL → np.array pour OCR
                    image_np = np.array(pil_image)
                    result = reader.readtext(image_np, detail=0)

                    if result:
                        md_lines.append(f"**Texte OCR (slide {i+1}, image {image_counter}) :**\n")
                        md_lines.append("> " + "\n> ".join(result) + "\n")  # blockquote style

                    image_counter += 1
                except Exception as e:
                    md_lines.append(f"[Erreur OCR sur image : {e}]\n")

    return "\n".join(md_lines)

# Utilisation
if __name__ == "__main__":
    markdown_output = pptx_to_markdown_text("tests/samples/test-transcription.pptx")
    print(markdown_output)
